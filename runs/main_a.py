import os, sys, logging
import numpy as np
import tensorflow as tf
from tensorflow import keras
from tensorflow.keras import models
import matplotlib.pyplot as plt
import seaborn as sns
import tf_keras_vis as vis
import sklearn.metrics
import pandas as pd
import argparse, time, json
import cv2
import pptx, re, datetime
from pptx.util import Inches
import skimage.transform

sys.path.append('/workspace/bitbucket/MRA')

parser = argparse.ArgumentParser()
main_config = parser.add_argument_group('network setting (must be provided)')

main_config.add_argument('--data_path', type=str, dest='data_path', default='/workspace/MRA')
main_config.add_argument('--data_type', type=str, dest='data_type', default='clinical')
main_config.add_argument('--excel_name', type=str, dest='excel_name', default='snubh01')
main_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp001')
main_config.add_argument('--train_name', type=str, dest='train_name', default='1,2,3,4,5')  # ex: 'train'
main_config.add_argument('--val_name', type=str, dest='val_name', default='6')  # ex: 'val'
main_config.add_argument('--model_name', type=str, dest='model_name', default='Model03')
main_config.add_argument('--trial_serial', type=int, dest='trial_serial', default=1)
main_config.add_argument('--image_size', type=int, dest='image_size', default=128)
main_config.add_argument('--channel_size', type=int, dest='channel_size', default=1)
main_config.add_argument('--seq_len', type=int, dest='seq_len', default=16)
main_config.add_argument('--num_classes', type=int, dest='num_classes', default=1)
main_config.add_argument('--max_keep', type=int, dest='max_keep', default=3)  # only use training
main_config.add_argument('--num_weight', type=int, dest='num_weight', default=1)  # only use validation
main_config.add_argument('--train', type=lambda x: x.title() in str(True), dest='train', default=False)
main_config.add_argument('--learning_rate', type=float, dest='learning_rate', default=0.001)
main_config.add_argument('--decay_steps', type=int, dest='decay_steps', default=5000)
main_config.add_argument('--decay_rate', type=int, dest='decay_rate', default=0.94)
main_config.add_argument('--batch_size', type=int, dest='batch_size', default=8)
main_config.add_argument('--epoch', type=int, dest='epoch', default=80)
main_config.add_argument('--alpha', type=float, dest='alpha', default=0.05)
main_config.add_argument('--gamma', type=float, dest='gamma', default=2.)
main_config.add_argument('--block_rep', type=str, dest='block_rep', default='4,6,6')
main_config.add_argument('--use_se', type=lambda x: x.title() in str(True), dest='use_se', default=False)
main_config.add_argument('--is_3d', type=lambda x: x.title() in str(True), dest='is_3d', default=False)
main_config.add_argument('--radius', type=int, dest='radius', default=80)


config, unparsed = parser.parse_known_args()

os.environ["TF_CPP_MIN_LOG_LEVEL"] = "3"
logging.disable(logging.WARNING)


# os.environ["CUDA_VISIBLE_DEVICES"] = "0"
gpu = tf.config.list_physical_devices('GPU')
tf.config.experimental.set_memory_growth(gpu[0], True)  # dynamic memory allocation
sns.set()  # apply seaborn style

trial_serial_str = '%03d' % config.trial_serial

log_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'logs-%s' % trial_serial_str)
result_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'result-%s' % trial_serial_str)
plot_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'plot-%s' % trial_serial_str)

if not os.path.exists(log_path): os.makedirs(log_path)
if not os.path.exists(result_path): os.makedirs(result_path)
if not os.path.exists(plot_path): os.makedirs(plot_path)


ROOT_PATH = '/workspace/MRA/'
RAW_PATH = os.path.join(ROOT_PATH, 'RAW', config.data_type)
INFO_PATH = os.path.join(ROOT_PATH, 'info')

df_path = os.path.join(INFO_PATH, config.excel_name) + '.xlsx'
df = pd.read_excel(df_path)

from data.setup_a import DataSettingV2
import models.model_a as model_a

ds = DataSettingV2(df=df, raw_path=RAW_PATH, image_size=config.image_size, batch_size=config.batch_size,
                   train_type=config.train_name, val_type=config.val_name,
                   seq_len=config.seq_len, train_seq_interval=config.seq_len, is_3d=config.is_3d)

ds.val = ds.val.batch(config.batch_size)
val_length = ds.val_length
print('val length: ', val_length)

if config.train:
    ds.train = ds.train.batch(config.batch_size)
    print('train length: ', ds.train_length)

infer_name = 'Inference' + config.model_name
infer = getattr(model_a, infer_name)(depth=config.seq_len, img_h=config.image_size, img_w=config.image_size,
                                     img_c=1, f_num=[64, 128, 256, 512], is_training=config.train)

# infer = getattr(model_a, infer_name)(input_size=[256, 256, 1], growth_k=32, image_size=config.image_size,
#                                      block_rep=config.block_rep, theta=0.5, use_se=config.use_se)
model = infer.model
cam_model = infer.cam_model

loss_fn = model_a.focal_loss_sigmoid
alpha, gamma = config.alpha, config.gamma
# loss_fn = model_a.focal_loss_sigmoid(alpha=config.alpha, gamma=config.gamma)


def training():
    info_log = {
        'EXCEL_FILE': config.excel_name,
        'TRAIN_NAME': config.train_name,
        'VAL_NAME': config.val_name,
        'IMAGE_SIZE': config.image_size,
        'BATCH_SIZE': config.batch_size,
        'LEARNING_RATE': config.learning_rate,
        'DECAY_STEPS': config.decay_steps,
        'DECAY_RATE': config.decay_rate,
        'EPOCH': config.epoch,
        'FOCAL_LOSS_ALPHA': config.alpha,
        'FOCAL_LOSS_GAMMA': config.gamma,
        'USE_SE': config.use_se,
        'BLOCK_REP': config.block_rep,
        'SEQ_LEN': config.seq_len,
        'IS_3D': config.is_3d
    }

    with open(os.path.join(result_path, '.info'), 'w') as f:
        f.write(json.dumps(info_log, indent=4, sort_keys=True))
        f.close()

    result_name = '_'.join([config.exp_name, config.model_name, trial_serial_str])+'.csv'
    auc_csv = pd.DataFrame({'WEIGHT_PATH': pd.Series(), 'AUC': pd.Series()})

    lr_schedule = keras.optimizers.schedules.ExponentialDecay(
        initial_learning_rate=config.learning_rate,
        decay_steps=config.decay_steps,
        decay_rate=config.decay_rate,
        staircase=True)

    optimizer = keras.optimizers.RMSprop(learning_rate=lr_schedule)

    # duration = 0.0
    perf_per_epoch, max_perf_per_epoch, max_current_step = [], [], []
    log_string = ''
    start_time = datetime.datetime.now()

    try:
        for epoch in range(1, config.epoch+1):
            train_loss, train_x, train_y = [], [], []
            for train_step, (img, lbl, name) in enumerate(ds.train):
                with tf.GradientTape() as tape:
                    # _, probs = model(img)
                    probs = model(img)
                    lbl = tf.cast(lbl, tf.float32)
                    train_loss_batch = loss_fn(lbl, probs, alpha, gamma)

                grads = tape.gradient(train_loss_batch, model.trainable_variables)
                optimizer.apply_gradients(zip(grads, model.trainable_variables))

                train_loss.append(train_loss_batch)
                train_x.extend(probs.numpy())
                train_y.extend(lbl.numpy())

                sys.stdout.write('Step: {0:>4d} ({1})\r'.format(train_step, epoch))

            try:
                fpr, tpr, _ = sklearn.metrics.roc_curve(train_y, train_x, drop_intermediate=False)
                train_auc = sklearn.metrics.auc(fpr, tpr)
            except:
                train_auc = 0.0000

            val_loss, val_x, val_y = [], [], []
            for val_step, (img, lbl, name) in enumerate(ds.val):
                sys.stdout.write('Evaluation [{0}/{1}]\r'.format(val_step, val_length // config.batch_size))

                # _, val_probs = model(img)
                val_probs = model(img)
                lbl = tf.cast(lbl, tf.float32)

                val_loss_batch = loss_fn(lbl, val_probs, alpha, gamma)
                val_loss.append(val_loss_batch)

                val_x.extend(val_probs.numpy())
                val_y.extend(lbl.numpy())

            try:
                fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, val_x, drop_intermediate=False)
                val_auc = sklearn.metrics.auc(fpr, tpr)
            except:
                val_auc = 0.0000

            time_elapsed = str(datetime.datetime.now() - start_time)
            log_string += 'Time Elapsed: {0}'.format(time_elapsed.split('.')[0])

            print('Epoch: %s Train-Loss: %.4f Train-AUC %.4f Val-Loss %.4f Val-AUC %.4f ' %
                  (epoch, np.mean(train_loss), train_auc, np.mean(val_loss), val_auc) + log_string)

            log_string = ''
            perf_per_epoch.append(val_auc)

            if epoch < config.max_keep + 1:
                max_current_step.append(epoch)
                max_perf_per_epoch.append(val_auc)

                weight_path = os.path.join(log_path, 'ckpt-' + '%03d' % epoch + '.hdf5')
                model.save(weight_path)

                auc_csv.loc[epoch] = weight_path, val_auc

            elif val_auc > min(auc_csv['AUC'].tolist()):
                os.remove(auc_csv.loc[max_current_step[0], 'WEIGHT_PATH'])
                auc_csv = auc_csv.drop(max_current_step[0])
                max_current_step.pop(0)
                max_current_step.append(epoch)
                max_perf_per_epoch.pop(0)
                max_perf_per_epoch.append(val_auc)

                weight_path = os.path.join(log_path, 'ckpt-' + '%03d' % epoch + '.hdf5')
                model.save(weight_path)

                auc_csv.loc[epoch] = weight_path, val_auc

            auc_csv.to_csv(os.path.join(result_path, result_name))

            if epoch == config.epoch: break

    except KeyboardInterrupt:
        print('Result saved')
        auc_csv.to_csv(os.path.join(result_path, result_name))


def gen_grad_cam(cam_layer, seq_len, loss, tape):
    grads = tape.gradient(loss, cam_layer)
    weights = tf.reduce_mean(grads, axis=(1, 2, 3))
    cam = np.zeros(cam_layer.shape[0:4], dtype=np.float32)

    batch_size = cam_layer.shape[0]

    heatmaps = np.zeros((batch_size, 256, 256, seq_len), dtype=np.float32)

    for batch in range(batch_size):  # batch size
        for index, w in enumerate(weights[batch]):  # each weights of batch
            cam[batch, :, :, :] += w * cam_layer[batch, :, :, :, index]
        cam_resize = skimage.transform.resize(cam[batch, :, :, :], (256, 256, seq_len))
        cam_resize = np.maximum(cam_resize, 0)  # ReLU
        heatmaps[batch, :, :, :] = (cam_resize - cam_resize.min()) / (cam_resize.max() - cam_resize.min())

    heatmaps = np.expand_dims(heatmaps, axis=-1)

    return heatmaps


def validation():
    weight_auc_path = os.path.join(config.data_path, config.exp_name, config.model_name,
                                   'result-%03d' % config.trial_serial)
    weight_auc_csv = pd.read_csv(os.path.join(weight_auc_path, '_'.join([config.exp_name, config.model_name,
                                                                         '%03d' % config.trial_serial])+'.csv'))
    weight_auc_csv = weight_auc_csv.sort_values('AUC', ascending=False)
    all_ckpt_paths = list(weight_auc_csv['WEIGHT_PATH'][0:int(config.num_weight)])

    imgs = np.zeros([ds.val_length, infer.img_h, infer.img_w, config.seq_len, infer.img_c])
    # cams = np.zeros([len(all_ckpt_paths), ds.val_length, infer.cam_h, infer.cam_w, infer.cam_d, infer.cam_f])
    cams = np.zeros([len(all_ckpt_paths), ds.val_length, infer.img_h, infer.img_w, config.seq_len, infer.img_c])

    lbls = np.zeros([ds.val_length, ], dtype=np.int32)
    probs = np.zeros([len(all_ckpt_paths), ds.val_length])
    names = []

    ckpt_idx = 0
    for ckpt in all_ckpt_paths:
        model.load_weights(ckpt)

        val_x, val_y, val_loss = [], [], []
        for step, (img, lbl, name) in enumerate(ds.val):
            sys.stdout.write('Evaluation [{0}/{1}]\r'.format(step, ds.val_length // config.batch_size))

            with tf.GradientTape() as tape:
                cam_layers, val_probs = cam_model(img)

            grad_cams = gen_grad_cam(cam_layers, config.seq_len, val_probs, tape)
            val_probs = np.squeeze(val_probs)
            lbl = tf.cast(lbl, tf.float32)
            loss_batch = loss_fn(lbl, val_probs, alpha, gamma)
            val_loss.append(loss_batch)

            val_x.extend(val_probs)
            val_y.extend(lbl.numpy())

            cams[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = grad_cams
            probs[ckpt_idx, step * config.batch_size:step * config.batch_size + len(lbl)] = val_probs

            if ckpt_idx == 0:
                imgs[step * config.batch_size:step * config.batch_size + len(lbl)] = img.numpy()
                lbls[step * config.batch_size:step * config.batch_size + len(lbl)] = lbl.numpy()
                names.extend(name.numpy().astype('str'))

        ckpt_idx += 1

    probs, cams = np.mean(probs, axis=0), np.mean(cams, axis=0)

    fpr, tpr, _ = sklearn.metrics.roc_curve(lbls, probs, drop_intermediate=False)
    val_auc = sklearn.metrics.auc(fpr, tpr)

    print('Validation AUC: %.4f' % val_auc)

    result_csv = pd.DataFrame({'NUMBER': names, 'PROB': np.squeeze(probs), 'LABEL': lbls})
    result_name = '_'.join([config.model_name, config.val_name, trial_serial_str, '%03d' % config.num_weight])+'.csv'
    result_csv.to_csv(os.path.join(result_path, result_name), index=False)

    prs = pptx.Presentation()
    prs.slide_width, prs.slide_height = Inches(8*2), Inches(4*2)

    plt_batch = 1
    plt_step = 0
    plt_iter, plt_examples = int(np.ceil(ds.val_length / plt_batch)), ds.val_length

    while plt_step < plt_iter:

        len_batch = plt_batch if plt_examples >= plt_batch else plt_examples

        images_batch = imgs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        labels_batch = lbls[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        names_batch = names[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        probs_batch = probs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        cams_batch = cams[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        show_seq_cam(cams_batch, probs_batch, images_batch, labels_batch, names_batch, 'LABEL')
        fig_name = '_'.join([config.exp_name, config.model_name, config.val_name, trial_serial_str,
                             '%03d' % plt_step]) + '.png'
        fig_path = os.path.join(plot_path, fig_name)
        plt.savefig(fig_path, bbox_inches='tight')
        slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(8 * 2), height=Inches(4 * 2))

        os.remove(fig_path)

        plt_step += 1
        plt_examples -= plt_batch

    print('plt_examples check: ', plt_examples)
    ppt_name = os.path.join(plot_path, '_'.join([config.exp_name, config.model_name, config.val_name, trial_serial_str,
                                                '%03d' % config.num_weight]) + '.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_seq_cam(cams, probs, images, labels, names, side_label, num_rows=4, num_cols=8, fig_size=(8*2, 4*2)):
    batch_size = cams.shape[0]
    fig, ax = plt.subplots(num_rows, num_cols, figsize=fig_size)
    axoff_fun = np.vectorize(lambda ax: ax.axis('off'))
    axoff_fun(ax)

    seq_len = config.seq_len
    for i in range(batch_size):
        prob = '%.2f' % probs[i]
        lbl = int(labels[i])

        ori_color = 'red' if lbl == 1 else 'blue'
        cam_color = 'red' if probs[i] >= 0.5 else 'blue'

        for j in range(seq_len):
            img_row, img_col = int(j % num_rows), int(j / num_rows) * 2

            show_image = np.squeeze(images[i, :, :, j, :])
            cam = np.squeeze(cams[i, :, :, j, :])

            ori_title = ' '.join([names[i], str(j), side_label + ': ' + str(lbl)])
            cam_title = side_label + ' Pred: ' + str(prob)

            ax[img_row, img_col].imshow(show_image, cmap='bone')
            ax[img_row, img_col + 1].imshow(show_image, cmap='bone')
            ax[img_row, img_col + 1].imshow(cam, cmap=plt.cm.seismic, alpha=0.5, interpolation='nearest')

            ax[img_row, img_col].set_title(ori_title, fontsize=7, color=ori_color)
            ax[img_row, img_col + 1].set_title(cam_title, fontsize=7, color=cam_color)


def show_cam(cams, probs, images, labels, names, side_label, num_rows=5, num_cols=8, fig_size=(8*2, 5*2)):
    batch_size = cams.shape[0]
    fig, ax = plt.subplots(num_rows, num_cols, figsize=fig_size)
    axoff_fun = np.vectorize(lambda ax: ax.axis('off'))
    axoff_fun(ax)

    for i in range(batch_size):
        prob = '%.2f' % probs[i]
        lbl = int(labels[i])
        show_image = np.squeeze(images[i])
        cam = np.squeeze(cams[i])
        img_row, img_col = int(i % num_rows), int(i / num_rows) * 2

        ori_title = ' '.join([names[i], side_label + ': '+str(lbl)])
        cam_title = side_label+' Pred: '+str(prob)

        ax[img_row, img_col].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(cam, cmap=plt.cm.seismic, alpha=0.5, interpolation='nearest')

        ori_color = 'red' if lbl == 1 else 'blue'
        cam_color = 'red' if probs[i] >= 0.5 else 'blue'

        ax[img_row, img_col].set_title(ori_title, fontsize=7, color=ori_color)
        ax[img_row, img_col+1].set_title(cam_title, fontsize=7, color=cam_color)


if __name__ == '__main__':
    if config.train:
        print('Training')
        training()
    else:
        print('Validation')
        validation()

